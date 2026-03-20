"""
Generate Astro Shelving Website Audit Presentation (PPTX)
Outputs to Desktop for easy sharing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0xFF, 0x6D, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
RED = RGBColor(0xE5, 0x39, 0x35)
AMBER = RGBColor(0xFF, 0xB3, 0x00)
SECTION_BG = RGBColor(0xF5, 0xF5, 0xF5)
DARK_CARD = RGBColor(0x2D, 0x2D, 0x44)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, height,
                              font_size=13, color=DARK_TEXT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BG)
add_shape(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "WEBSITE AUDIT REPORT", 18, LIGHT_GRAY, font_name="Calibri")
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
             "Astro Shelving", 52, WHITE, True, font_name="Calibri")
add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "Lead & Sales Generation Analysis", 28, ACCENT, font_name="Calibri")
add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "astroshelving.co.za", 18, LIGHT_GRAY, font_name="Calibri")

add_text_box(slide, Inches(1), Inches(6.2), Inches(5), Inches(0.5),
             "Prepared by AnyVision Media  |  March 2026", 14, LIGHT_GRAY, font_name="Calibri")
add_text_box(slide, Inches(7), Inches(6.2), Inches(5.5), Inches(0.5),
             "ian@anyvisionmedia.com", 14, ACCENT, alignment=PP_ALIGN.RIGHT, font_name="Calibri")

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "AGENDA", 32, WHITE, True)

items_left = [
    "1.  Executive Summary & Health Score",
    "2.  Critical Weaknesses (Must Fix)",
    "3.  Technical SEO Issues",
    "4.  Content & Conversion Gaps",
    "5.  Missing Pages That Drive Leads",
    "6.  Local SEO Opportunities",
]
items_right = [
    "7.   Keyword Opportunities",
    "8.   Schema & Structured Data",
    "9.   Competitive Advantages",
    "10.  Conversion Optimization Roadmap",
    "11.  Estimated Impact & ROI",
    "12.  Next Steps",
]
add_bullet_slide_content(slide, items_left, Inches(1.5), Inches(2), Inches(5), Inches(4.5), 18, DARK_TEXT, Pt(14))
add_bullet_slide_content(slide, items_right, Inches(7), Inches(2), Inches(5), Inches(4.5), 18, DARK_TEXT, Pt(14))

# ============================================================
# SLIDE 3: EXECUTIVE SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "EXECUTIVE SUMMARY", 32, WHITE, True)

# Score circle area
add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.2), Inches(4.8), SECTION_BG)
add_text_box(slide, Inches(1), Inches(2.0), Inches(2.8), Inches(0.5),
             "OVERALL HEALTH SCORE", 14, DARK_TEXT, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(2.8), Inches(2.8), Inches(1.5),
             "38", 72, RED, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(2.8), Inches(0.5),
             "out of 100", 18, DARK_TEXT, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.8), Inches(2.8), Inches(0.5),
             "NEEDS SIGNIFICANT WORK", 13, RED, True, PP_ALIGN.CENTER)

# Key stats
stats = [
    ("Website", "astroshelving.co.za"),
    ("Business", "Astro Storage Solutions (Pty) Ltd"),
    ("Founded", "November 2009 (15+ years)"),
    ("Location", "Edenvale, Johannesburg"),
    ("Products", "Shelving, Racking, Mezzanines"),
    ("Pages Indexed", "39"),
    ("Pages with Pricing", "1 of 39"),
    ("Lead Capture Methods", "Contact form only"),
    ("Blog Posts", "0"),
    ("Case Studies", "0"),
]
y = Inches(1.8)
for label, value in stats:
    add_text_box(slide, Inches(4.5), y, Inches(2.5), Inches(0.35), label + ":", 13, DARK_TEXT, True)
    add_text_box(slide, Inches(7), y, Inches(5.5), Inches(0.35), value, 13, DARK_TEXT)
    y += Inches(0.4)

# Bottom insight
add_shape(slide, Inches(4.5), Inches(5.8), Inches(8), Inches(1.2), DARK_CARD)
add_text_box(slide, Inches(4.8), Inches(5.95), Inches(7.5), Inches(0.9),
             "The website functions as a digital brochure, not a lead generation machine.\n"
             "B2B buyers complete 70% of research BEFORE contacting a supplier.\n"
             "Every missing piece of information is a lead lost to a competitor.",
             12, WHITE)

# ============================================================
# SLIDE 4: SCORE BREAKDOWN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "SCORE BREAKDOWN BY CATEGORY", 32, WHITE, True)

categories = [
    ("Technical SEO", "45/100", AMBER, [
        "HTTPS: Yes (good)",
        "Sitemap exists but contains junk pages",
        "Only basic Organization schema",
        "Conflicting business hours (schema vs contact page)",
        "Most pages last updated Feb 2024 (stale)"
    ]),
    ("Content Quality", "25/100", RED, [
        "Thin product descriptions (200-400 words)",
        "No blog, guides, or case studies",
        "Location pages: 60% duplicate content",
        "\"Other Services\" page lists items with zero detail"
    ]),
    ("Conversion", "20/100", RED, [
        "Single conversion path (contact form + phone)",
        "No WhatsApp, live chat, or quote calculator",
        "No pricing transparency on any page",
        "No lead magnets or email capture"
    ]),
    ("Trust & Authority", "35/100", AMBER, [
        "Only 3 testimonials (no photos/company names)",
        "No project gallery or case studies",
        "No certifications mentioned",
        "Copyright shows 2022"
    ]),
    ("Local SEO", "50/100", AMBER, [
        "16 location pages (but thin/doorway quality)",
        "Good address/phone visibility",
        "Missing major areas (Sandton, Pretoria, Midrand)",
        "No Google Business Profile optimization"
    ]),
]

x_start = Inches(0.5)
card_width = Inches(2.4)
card_gap = Inches(0.15)
for i, (name, score, color, items) in enumerate(categories):
    x = x_start + i * (card_width + card_gap)
    # Card background
    add_shape(slide, x, Inches(1.6), card_width, Inches(5.5), SECTION_BG)
    # Score header
    add_shape(slide, x, Inches(1.6), card_width, Inches(1.2), DARK_CARD)
    add_text_box(slide, x, Inches(1.7), card_width, Inches(0.4),
                 name, 13, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.1), card_width, Inches(0.6),
                 score, 28, color, True, PP_ALIGN.CENTER)
    # Items
    add_bullet_slide_content(slide, [f"  {item}" for item in items],
                              x + Inches(0.1), Inches(3.0), card_width - Inches(0.2),
                              Inches(4), 10, DARK_TEXT, Pt(4))

# ============================================================
# SLIDE 5: CRITICAL #1 - NO LEAD CAPTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x8B, 0x00, 0x00))
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
             "CRITICAL WEAKNESS #1", 16, RGBColor(0xFF, 0xCC, 0xCC))
add_text_box(slide, Inches(0.8), Inches(0.55), Inches(11), Inches(0.6),
             "No Lead Capture System", 30, WHITE, True)

# Current state
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(3.5), SECTION_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4),
             "CURRENT STATE", 16, RED, True)
problems = [
    "Only ONE conversion path: basic contact form",
    "No WhatsApp click-to-chat (essential in SA)",
    "No email capture (newsletters, downloads)",
    "No lead magnets (guides, calculators)",
    "No exit-intent popups",
    "No sticky CTAs on any page",
]
add_bullet_slide_content(slide, [f"  X  {p}" for p in problems],
                          Inches(1), Inches(2.2), Inches(5), Inches(2.8), 13, DARK_TEXT, Pt(6))

# Recommendations
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.5), SECTION_BG)
add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
             "RECOMMENDATIONS", 16, GREEN, True)
fixes = [
    "Add WhatsApp floating button on every page",
    "Add sticky \"Get a Free Quote\" bar site-wide",
    "Build multi-step form: Type > Size > Qty > Info",
    "Create lead magnet: Free Storage Planning Guide",
    "Add exit-intent popup for guide download",
]
add_bullet_slide_content(slide, [f"  ->  {f}" for f in fixes],
                          Inches(7), Inches(2.2), Inches(5.5), Inches(2.8), 13, DARK_TEXT, Pt(6))

# Impact bar
add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.2), DARK_CARD)
add_text_box(slide, Inches(1), Inches(5.65), Inches(11.5), Inches(0.9),
             "IMPACT: Losing 70-80% of potential leads who aren't ready to call or email.\n"
             "ESTIMATED IMPROVEMENT: +20-30% more leads from WhatsApp alone.",
             14, WHITE)

# ============================================================
# SLIDE 6: CRITICAL #2 - NO PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x8B, 0x00, 0x00))
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
             "CRITICAL WEAKNESS #2", 16, RGBColor(0xFF, 0xCC, 0xCC))
add_text_box(slide, Inches(0.8), Inches(0.55), Inches(11), Inches(0.6),
             "Zero Pricing Transparency", 30, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.8), SECTION_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4),
             "CURRENT STATE", 16, RED, True)
add_bullet_slide_content(slide, [
    "  X  No prices on ANY product page",
    "  X  Only 1 exception: R850 offer page",
    "  X  Only 1 product in WooCommerce (from 2021)",
    "  X  Buyers research pricing online BEFORE calling",
], Inches(1), Inches(2.2), Inches(5), Inches(2), 13, DARK_TEXT, Pt(6))

add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.8), SECTION_BG)
add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
             "RECOMMENDATIONS", 16, GREEN, True)
add_bullet_slide_content(slide, [
    "  ->  Add \"Starting from R___\" on every product",
    "  ->  Create dedicated Pricing page with ranges",
    "  ->  Show typical project costs by category:",
    "       Bolt-on Shelving: From R850/unit",
    "       Pallet Racking: From R___/bay",
    "       Mezzanine: From R___/sqm",
    "  ->  Add \"Request Exact Quote\" CTA next to prices",
], Inches(7), Inches(2.2), Inches(5.5), Inches(2.5), 13, DARK_TEXT, Pt(4))

add_shape(slide, Inches(0.8), Inches(4.8), Inches(11.8), Inches(1), DARK_CARD)
add_text_box(slide, Inches(1), Inches(4.9), Inches(11.5), Inches(0.7),
             "ESTIMATED IMPROVEMENT: +15-25% more leads with pricing transparency", 14, WHITE)

# ============================================================
# SLIDE 7: CRITICAL #3 - THIN CONTENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x8B, 0x00, 0x00))
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
             "CRITICAL WEAKNESS #3", 16, RGBColor(0xFF, 0xCC, 0xCC))
add_text_box(slide, Inches(0.8), Inches(0.55), Inches(11), Inches(0.6),
             "Extremely Thin Content", 30, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), SECTION_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4),
             "CURRENT STATE", 16, RED, True)
add_bullet_slide_content(slide, [
    "  X  Product pages: 200-400 words (generic)",
    "  X  No buying guides or comparison content",
    "  X  No FAQ page",
    "  X  No blog or knowledge base",
    "  X  \"Other Services\" page: zero detail",
    "  X  Location pages: 60% duplicate content",
], Inches(1), Inches(2.2), Inches(5), Inches(2), 13, DARK_TEXT, Pt(6))

add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2), SECTION_BG)
add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
             "WHAT'S NEEDED", 16, GREEN, True)
add_bullet_slide_content(slide, [
    "  ->  Each product page: 800-1,500 words",
    "       - Specifications tables",
    "       - Load capacity charts",
    "       - Use cases (warehouse, retail, office)",
    "       - Installation photos (before/after)",
    "       - Comparison vs alternatives",
    "",
    "  ->  Blog content (2-4 posts/month):",
    "       - \"5 Signs Your Warehouse Needs Better Shelving\"",
    "       - \"Warehouse Safety Regulations in SA\"",
    "       - \"How to Maximize Storage Space\"",
    "       - \"Pallet Racking Load Capacity Guide\"",
], Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), 12, DARK_TEXT, Pt(3))

# ============================================================
# SLIDE 8: CRITICAL #4 - WEAK TRUST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), RGBColor(0x8B, 0x00, 0x00))
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
             "CRITICAL WEAKNESS #4", 16, RGBColor(0xFF, 0xCC, 0xCC))
add_text_box(slide, Inches(0.8), Inches(0.55), Inches(11), Inches(0.6),
             "Weak Trust Signals", 30, WHITE, True)

add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4), SECTION_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4),
             "CURRENT STATE", 16, RED, True)
add_bullet_slide_content(slide, [
    "  X  Only 3 testimonials (no photos, no companies)",
    "  X  No case studies or project portfolio",
    "  X  No certifications (SABS, ISO, etc.)",
    "  X  No client logos or \"projects completed\" count",
    "  X  Copyright shows 2022 (outdated)",
    "  X  No Google Reviews integration",
    "  X  No video content",
], Inches(1), Inches(2.2), Inches(5), Inches(3), 13, DARK_TEXT, Pt(6))

add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4), SECTION_BG)
add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
             "RECOMMENDATIONS", 16, GREEN, True)
add_bullet_slide_content(slide, [
    "  ->  Projects/Gallery page: 20+ real photos",
    "  ->  3-5 case studies with client quotes",
    "  ->  Embed Google Reviews widget",
    "  ->  Add \"Since 2009\" trust badge to header",
    "  ->  Counter bar: 500+ Projects | 15+ Years",
    "  ->  Display certifications and compliance",
    "  ->  Update copyright to 2026",
    "  ->  Create 2-3 video testimonials",
], Inches(7), Inches(2.2), Inches(5.5), Inches(3), 13, DARK_TEXT, Pt(6))

# ============================================================
# SLIDE 9: TECHNICAL SEO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "TECHNICAL SEO ISSUES", 32, WHITE, True)

issues = [
    ("Poor Title Tags", "Homepage: \"Home - Astro Shelving\" (generic, low CTR)",
     "Fix: \"Industrial Shelving & Racking | Johannesburg | Astro Shelving\""),
    ("Sitemap Junk Pages", "/sample-page/ (2019), /cart/, /checkout/, /my-account/ (empty WooCommerce)",
     "Fix: Remove or noindex these pages immediately"),
    ("Conflicting Hours", "Schema: Mon-Sun 09:00-17:00 vs Contact: Mon-Fri 08:00-16:30, Sat 07:00-14:00",
     "Fix: Align both to actual business hours"),
    ("Minimal Schema", "Only basic Organization schema. Missing Product, BreadcrumbList, LocalBusiness",
     "Fix: Add all recommended schemas (code provided)"),
    ("Stale Content", "Most pages last modified Feb 2024 (2+ years ago)",
     "Fix: Update and expand content quarterly minimum"),
]

y = Inches(1.6)
for title, problem, fix in issues:
    add_shape(slide, Inches(0.8), y, Inches(11.8), Inches(1), SECTION_BG)
    add_text_box(slide, Inches(1), y + Inches(0.05), Inches(2.5), Inches(0.35),
                 title, 14, RED, True)
    add_text_box(slide, Inches(3.6), y + Inches(0.05), Inches(4.5), Inches(0.4),
                 problem, 11, DARK_TEXT)
    add_text_box(slide, Inches(8.2), y + Inches(0.05), Inches(4.2), Inches(0.4),
                 fix, 11, GREEN)
    y += Inches(1.1)

# ============================================================
# SLIDE 10: MISSING PAGES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "MISSING PAGES THAT WOULD DRIVE LEADS", 32, WHITE, True)

# High priority
add_shape(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(3.5), Inches(0.4),
             "HIGH PRIORITY", 16, RED, True)
add_bullet_slide_content(slide, [
    "/pricing/",
    "  How Much Does Shelving Cost?",
    "  Captures highest-intent searches",
    "",
    "/projects/ or /gallery/",
    "  Photo portfolio of completed work",
    "  Builds trust, shows capability",
    "",
    "/faqs/",
    "  15-20 common questions",
    "  Captures informational searches",
    "",
    "/blog/",
    "  Content hub (2-4 posts/month)",
    "  Drives organic traffic over time",
], Inches(0.7), Inches(2.1), Inches(3.5), Inches(4.5), 11, DARK_TEXT, Pt(2))

# Medium priority
add_shape(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_text_box(slide, Inches(4.9), Inches(1.6), Inches(3.5), Inches(0.4),
             "MEDIUM PRIORITY", 16, AMBER, True)
add_bullet_slide_content(slide, [
    "/industries/warehouse/",
    "  Industry-specific landing page",
    "",
    "/industries/retail/",
    "  Retail shop fitting focus",
    "",
    "/industries/office/",
    "  Office storage solutions",
    "",
    "/compare/shelving-vs-racking/",
    "  Comparison content",
    "",
    "/buying-guide/",
    "  How to Choose Right Shelving",
], Inches(4.9), Inches(2.1), Inches(3.5), Inches(4.5), 11, DARK_TEXT, Pt(2))

# Low priority
add_shape(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_text_box(slide, Inches(9.1), Inches(1.6), Inches(3.5), Inches(0.4),
             "LOW PRIORITY", 16, GREEN, True)
add_bullet_slide_content(slide, [
    "/storage-calculator/",
    "  Interactive lead magnet tool",
    "  Users input dimensions -> get estimate",
    "",
    "/delivery-areas/",
    "  Expand geo-targeting beyond East Rand",
    "  Cover Sandton, Pretoria, Midrand",
    "",
    "/careers/",
    "  Shows company growth & culture",
    "",
    "/case-studies/",
    "  Detailed project breakdowns",
    "  Problem -> Solution -> Results",
], Inches(9.1), Inches(2.1), Inches(3.5), Inches(4.5), 11, DARK_TEXT, Pt(2))

# ============================================================
# SLIDE 11: LOCAL SEO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "LOCAL SEO OPPORTUNITIES", 32, WHITE, True)

# Current
add_shape(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), SECTION_BG)
add_text_box(slide, Inches(1), Inches(1.7), Inches(5), Inches(0.4),
             "16 EXISTING LOCATION PAGES", 14, DARK_TEXT, True)
add_bullet_slide_content(slide, [
    "Edenvale, Bedfordview, Kempton Park, Spartan,",
    "Isando, Jet Park, Wilbart, Sunnyrock,",
    "Modderfontein, Eastleigh (+ mezzanine variants)",
    "",
    "PROBLEM: 60% duplicate/boilerplate content",
    "These are semi-thin doorway pages"
], Inches(1), Inches(2.2), Inches(5), Inches(2), 12, DARK_TEXT, Pt(4))

# Missing areas
add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5), SECTION_BG)
add_text_box(slide, Inches(7), Inches(1.7), Inches(5.5), Inches(0.4),
             "MISSING HIGH-VALUE AREAS", 14, RED, True)
add_bullet_slide_content(slide, [
    "Sandton / Randburg / Fourways",
    "Pretoria / Centurion / Midrand",
    "Boksburg / Benoni / Springs",
    "Germiston / Alberton",
    "Cape Town / Durban (if national)",
], Inches(7), Inches(2.2), Inches(5.5), Inches(2), 12, DARK_TEXT, Pt(6))

# Action plan
add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.8), Inches(2.5), DARK_CARD)
add_text_box(slide, Inches(1), Inches(4.6), Inches(11.5), Inches(0.4),
             "ACTION PLAN", 16, ACCENT, True)
add_bullet_slide_content(slide, [
    "1.  Rewrite existing 16 pages with unique content (500+ words each, area-specific photos)",
    "2.  Create 8-10 new location pages for high-value areas listed above",
    "3.  Optimize Google Business Profile (photos, posts, Q&A, reviews)",
    "4.  Build local citations (Yell.co.za, Brabys, Yellow Pages SA, Cylex)",
    "5.  Launch review generation campaign (ask clients post-installation)",
], Inches(1), Inches(5.1), Inches(11.5), Inches(1.8), 13, WHITE, Pt(6))

# ============================================================
# SLIDE 12: KEYWORD OPPORTUNITIES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "KEYWORD OPPORTUNITIES", 32, WHITE, True)

# High intent table
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(6), Inches(0.4),
             "HIGH-INTENT KEYWORDS (Ready to Buy)", 16, RED, True)
hi_keywords = [
    ("industrial shelving south africa", "500-1,000/mo", "Medium"),
    ("pallet racking johannesburg", "200-500/mo", "Medium"),
    ("warehouse shelving for sale", "200-500/mo", "High"),
    ("shelving prices south africa", "100-300/mo", "Low"),
    ("shop shelving south africa", "200-500/mo", "Medium"),
    ("mezzanine floor johannesburg", "100-200/mo", "Low"),
    ("used shelving for sale jhb", "100-200/mo", "Low"),
    ("cantilever racking SA", "50-100/mo", "Low"),
]
y = Inches(2.0)
# Header row
add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.35), DARK_CARD)
add_text_box(slide, Inches(0.9), y, Inches(2.8), Inches(0.35), "Keyword", 11, WHITE, True)
add_text_box(slide, Inches(3.8), y, Inches(1.3), Inches(0.35), "Est. Volume", 11, WHITE, True)
add_text_box(slide, Inches(5.1), y, Inches(1.1), Inches(0.35), "Competition", 11, WHITE, True)
y += Inches(0.35)
for kw, vol, comp in hi_keywords:
    color = SECTION_BG if hi_keywords.index((kw, vol, comp)) % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y, Inches(5.5), Inches(0.3), color)
    add_text_box(slide, Inches(0.9), y, Inches(2.8), Inches(0.3), kw, 10, DARK_TEXT)
    add_text_box(slide, Inches(3.8), y, Inches(1.3), Inches(0.3), vol, 10, DARK_TEXT)
    comp_color = GREEN if comp == "Low" else (AMBER if comp == "Medium" else RED)
    add_text_box(slide, Inches(5.1), y, Inches(1.1), Inches(0.3), comp, 10, comp_color, True)
    y += Inches(0.3)

# Long tail
add_text_box(slide, Inches(7), Inches(1.5), Inches(6), Inches(0.4),
             "LONG-TAIL KEYWORDS (Low Competition)", 16, GREEN, True)
long_tail = [
    "how much does industrial shelving cost in SA",
    "best shelving for small warehouse",
    "drive-in racking vs selective racking",
    "warehouse shelving installation johannesburg",
    "custom mezzanine floor edenvale",
    "gondola shelving for retail shops",
    "document storage racking systems",
    "cold room shelving south africa",
    "heavy duty garage shelving johannesburg",
    "warehouse safety shelving regulations SA",
]
add_bullet_slide_content(slide, [f"  ->  {kw}" for kw in long_tail],
                          Inches(7), Inches(2.0), Inches(6), Inches(4.5), 12, DARK_TEXT, Pt(6))

# ============================================================
# SLIDE 13: SCHEMA & STRUCTURED DATA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "SCHEMA & STRUCTURED DATA GAPS", 32, WHITE, True)

schemas = [
    ("LocalBusiness", "Homepage", "Yes - Knowledge Panel", "Enables Google Knowledge Panel with business info, hours, location"),
    ("Product", "Each product page", "Yes - Product results", "Shows pricing, availability, ratings in search results"),
    ("BreadcrumbList", "All pages", "Yes - Breadcrumb trail", "Displays navigation path in SERPs, improves CTR"),
    ("WebSite + Search", "Homepage", "Yes - Sitelinks", "Enables sitelinks search box in Google"),
    ("Service", "Service pages", "Partial", "Describes installation, design, consultation services"),
    ("ContactPage", "Contact page", "No (AI search)", "Helps AI search engines understand contact info"),
    ("OpenGraph", "All pages", "Social previews", "Controls how links appear on Facebook, LinkedIn, WhatsApp"),
]

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.4),
             "Currently: Only basic Organization schema", 13, RED, True)

y = Inches(1.9)
# Header
add_shape(slide, Inches(0.8), y, Inches(11.8), Inches(0.4), DARK_CARD)
add_text_box(slide, Inches(0.9), y, Inches(2), Inches(0.4), "Schema Type", 12, WHITE, True)
add_text_box(slide, Inches(3), y, Inches(2), Inches(0.4), "Where", 12, WHITE, True)
add_text_box(slide, Inches(5.1), y, Inches(2.2), Inches(0.4), "Rich Result?", 12, WHITE, True)
add_text_box(slide, Inches(7.4), y, Inches(5), Inches(0.4), "Benefit", 12, WHITE, True)
y += Inches(0.4)

for schema, where, rich, benefit in schemas:
    bg = SECTION_BG if schemas.index((schema, where, rich, benefit)) % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y, Inches(11.8), Inches(0.45), bg)
    add_text_box(slide, Inches(0.9), y, Inches(2), Inches(0.45), schema, 11, DARK_TEXT, True)
    add_text_box(slide, Inches(3), y, Inches(2), Inches(0.45), where, 11, DARK_TEXT)
    add_text_box(slide, Inches(5.1), y, Inches(2.2), Inches(0.45), rich, 11, GREEN if "Yes" in rich else AMBER)
    add_text_box(slide, Inches(7.4), y, Inches(5), Inches(0.45), benefit, 10, DARK_TEXT)
    y += Inches(0.45)

add_shape(slide, Inches(0.8), y + Inches(0.3), Inches(11.8), Inches(0.6), DARK_CARD)
add_text_box(slide, Inches(1), y + Inches(0.35), Inches(11.5), Inches(0.5),
             "NOTE: FAQ schema no longer generates Google rich results (restricted Aug 2023). "
             "HowTo schema also deprecated (Sep 2023). All recommended code is ready to implement.",
             12, LIGHT_GRAY)

# ============================================================
# SLIDE 14: COMPETITIVE ADVANTAGES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "COMPETITIVE ADVANTAGES TO AMPLIFY", 32, WHITE, True)

advantages = [
    ("15+ Years", "Est. 2009", "Add \"Since 2009\" badge to header.\n\"Trusted for 15+ years\" in hero."),
    ("Free On-Site\nConsultation", "Major\ndifferentiator", "Deserves its own landing page.\n\"We come to you\" messaging."),
    ("Free Design\nService", "Lead\nmagnet", "\"Free 3D Warehouse Layout Design\"\nwith visual examples."),
    ("In-House\nManufacturing", "Quality\ncontrol", "\"Made in South Africa\" badge.\nShow factory + process photos."),
    ("Full Service\nPipeline", "End-to-end\nsolution", "Consult > Design > Manufacture >\nInstall > After-Sales visual."),
    ("Strategic\nLocation", "East Rand\nindustrial", "Fast delivery for Gauteng.\nProximity to industrial clients."),
]

x = Inches(0.5)
card_w = Inches(2)
for i, (title, subtitle, desc) in enumerate(advantages):
    cx = x + i * (card_w + Inches(0.13))
    add_shape(slide, cx, Inches(1.6), card_w, Inches(5.2), SECTION_BG)
    add_shape(slide, cx, Inches(1.6), card_w, Inches(1.5), ACCENT)
    add_text_box(slide, cx, Inches(1.7), card_w, Inches(0.8),
                 title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, cx, Inches(2.5), card_w, Inches(0.5),
                 subtitle, 11, WHITE, False, PP_ALIGN.CENTER)
    add_text_box(slide, cx + Inches(0.15), Inches(3.4), card_w - Inches(0.3), Inches(3),
                 desc, 11, DARK_TEXT)

# ============================================================
# SLIDE 15: ROADMAP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "CONVERSION OPTIMIZATION ROADMAP", 32, WHITE, True)

# Week 1-2
add_shape(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), GREEN)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(3.5), Inches(0.4),
             "WEEK 1-2: QUICK WINS", 14, WHITE, True)
add_bullet_slide_content(slide, [
    "Add WhatsApp click-to-chat button",
    "Add sticky \"Get Free Quote\" CTA bar",
    "Fix copyright to 2026",
    "Fix conflicting business hours",
    "Remove junk pages from sitemap",
    "Update meta titles for better CTR",
    "Add Google Reviews widget",
], Inches(0.7), Inches(2.2), Inches(3.5), Inches(4.5), 12, DARK_TEXT, Pt(8))

# Month 1-2
add_shape(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_shape(slide, Inches(4.7), Inches(1.5), Inches(4), Inches(0.5), AMBER)
add_text_box(slide, Inches(4.9), Inches(1.55), Inches(3.5), Inches(0.4),
             "MONTH 1-2: MEDIUM TERM", 14, WHITE, True)
add_bullet_slide_content(slide, [
    "Build multi-step quote request form",
    "Create FAQ page (15-20 questions)",
    "Create project gallery (20+ photos)",
    "Add pricing indicators on products",
    "Expand product pages to 800+ words",
    "Add Schema markup (code ready)",
    "Build 5 industry landing pages",
    "Google Business Profile optimization",
], Inches(4.9), Inches(2.2), Inches(3.5), Inches(4.5), 12, DARK_TEXT, Pt(6))

# Month 3-6
add_shape(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(5.5), SECTION_BG)
add_shape(slide, Inches(8.9), Inches(1.5), Inches(4), Inches(0.5), ACCENT)
add_text_box(slide, Inches(9.1), Inches(1.55), Inches(3.5), Inches(0.4),
             "MONTH 3-6: LONG TERM", 14, WHITE, True)
add_bullet_slide_content(slide, [
    "Launch blog (2 posts/month)",
    "Build storage calculator tool",
    "Create comparison/buying guides",
    "Expand location pages (new areas)",
    "Rewrite thin location pages",
    "Email nurture for unconverted leads",
    "Add retargeting pixels (Meta+Google)",
    "Review generation campaign",
], Inches(9.1), Inches(2.2), Inches(3.5), Inches(4.5), 12, DARK_TEXT, Pt(6))

# ============================================================
# SLIDE 16: ESTIMATED IMPACT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "ESTIMATED IMPACT & ROI", 32, WHITE, True)

impacts = [
    ("WhatsApp + Sticky CTA", "+20-30%"),
    ("Pricing Transparency", "+15-25%"),
    ("Multi-step Quote Form", "+10-20%"),
    ("FAQ Page + Schema", "+10-15%"),
    ("Project Gallery + Case Studies", "+10-15%"),
    ("Blog Content (6 months)", "+20-40%"),
    ("Local SEO Optimization", "+15-25%"),
]

y = Inches(1.8)
add_shape(slide, Inches(1.5), y, Inches(7), Inches(0.4), DARK_CARD)
add_text_box(slide, Inches(1.6), y, Inches(4), Inches(0.4), "Improvement", 13, WHITE, True)
add_text_box(slide, Inches(5.8), y, Inches(2.5), Inches(0.4), "Est. Lead Increase", 13, WHITE, True)
y += Inches(0.4)

for improvement, impact in impacts:
    bg = SECTION_BG if impacts.index((improvement, impact)) % 2 == 0 else WHITE
    add_shape(slide, Inches(1.5), y, Inches(7), Inches(0.4), bg)
    add_text_box(slide, Inches(1.6), y, Inches(4), Inches(0.4), improvement, 13, DARK_TEXT)
    add_text_box(slide, Inches(5.8), y, Inches(2.5), Inches(0.4), impact, 13, GREEN, True)
    y += Inches(0.4)

# Combined
add_shape(slide, Inches(1.5), y + Inches(0.1), Inches(7), Inches(0.6), ACCENT)
add_text_box(slide, Inches(1.6), y + Inches(0.15), Inches(4), Inches(0.5),
             "COMBINED 6-MONTH PROJECTION", 15, WHITE, True)
add_text_box(slide, Inches(5.8), y + Inches(0.15), Inches(2.5), Inches(0.5),
             "+80-150%", 22, WHITE, True)

# Key insight
add_shape(slide, Inches(1.5), y + Inches(1.2), Inches(7), Inches(1.5), DARK_CARD)
add_text_box(slide, Inches(1.7), y + Inches(1.3), Inches(6.5), Inches(1.3),
             "KEY INSIGHTS:\n"
             "- Every 1s of load time improvement = ~7% conversion increase\n"
             "- B2B buyers complete 70% of research before contacting suppliers\n"
             "- Current baseline is so low that even quick wins could double leads in 30 days",
             12, WHITE)

# ============================================================
# SLIDE 17: NEXT STEPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), DARK_BG)
add_shape(slide, Inches(0), Inches(5.8), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "NEXT STEPS", 36, WHITE, True)

# Immediate
add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2), DARK_CARD)
add_text_box(slide, Inches(1), Inches(1.55), Inches(5), Inches(0.4),
             "IMMEDIATE (This Week)", 16, ACCENT, True)
add_bullet_slide_content(slide, [
    "1.  Approve quick wins for implementation",
    "2.  Gather 20+ project photos",
    "3.  Collect 5-10 client testimonials",
    "4.  Confirm actual business hours",
], Inches(1), Inches(2.0), Inches(5), Inches(1.3), 13, WHITE, Pt(6))

# Short term
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2), DARK_CARD)
add_text_box(slide, Inches(7), Inches(1.55), Inches(5.5), Inches(0.4),
             "SHORT TERM (30 Days)", 16, ACCENT, True)
add_bullet_slide_content(slide, [
    "5.  Implement Schema markup (code ready)",
    "6.  Build FAQ page and project gallery",
    "7.  Add pricing ranges to products",
    "8.  Optimize Google Business Profile",
], Inches(7), Inches(2.0), Inches(5.5), Inches(1.3), 13, WHITE, Pt(6))

# Medium term
add_shape(slide, Inches(0.8), Inches(3.8), Inches(11.8), Inches(1.6), DARK_CARD)
add_text_box(slide, Inches(1), Inches(3.85), Inches(11.5), Inches(0.4),
             "MEDIUM TERM (60-90 Days)", 16, ACCENT, True)
add_bullet_slide_content(slide, [
    "9.  Expand product page content (800+ words each)          "
    "10.  Launch blog with first 4 posts          "
    "11.  Build industry-specific landing pages          "
    "12.  Set up review generation campaign",
], Inches(1), Inches(4.3), Inches(11.5), Inches(0.8), 13, WHITE, Pt(4))

# Contact
add_text_box(slide, Inches(1), Inches(6.0), Inches(5), Inches(0.5),
             "Prepared by AnyVision Media", 16, LIGHT_GRAY)
add_text_box(slide, Inches(1), Inches(6.4), Inches(5), Inches(0.5),
             "ian@anyvisionmedia.com", 14, ACCENT)
add_text_box(slide, Inches(7), Inches(6.2), Inches(5.5), Inches(0.5),
             "March 2026", 14, LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# ============================================================
# SAVE
# ============================================================
output_path = os.path.expanduser("~/OneDrive/Desktop/Astro_Shelving_Website_Audit.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
